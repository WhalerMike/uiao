"""Canonical AAN (Application-Aware Networking) PowerPoint briefing-deck generator.

Reproduces the *Book_10 / Book_11 OrgComp* deck DNA — native-shape composition
on a white 10.0 x 5.625 in (16:9) canvas, Arial typography, the exact UIAO
house palette, full presenter notes on every slide, and the progressive
build-sequence reveal with the amber-ring active step. The authoritative
measurements this generator reproduces live in
``docs/customer-documents/orgcomp-series/OrgComp_DECK_STYLE_NOTES.md``.

Design intent
-------------
* **Data-driven** — decks are authored as a YAML *deck spec* (see the schema
  below). No hand-editing of the ``.pptx``; the ``.pptx`` is a build artifact.
* **White only** — every slide gets a solid white background; there are no dark
  title or closing slides (see the "No dark slides" rule in the style notes).
* **Sourced** — vendor claims append a single ``Vendor sources: ...`` line to the
  speaker notes, with URLs resolved against the :data:`SRC` link-map (a deck spec
  may extend the map via its ``sources:`` key).
* **Date Code** — the caller passes ``date_code`` explicitly (do NOT derive it
  from ``TZ`` on this machine, which returns UTC). Format
  ``YYYY-MM-DD HH:MM ET``; it is rendered verbatim on the title slide.

Deck-spec YAML schema
=====================
Top level has two keys: ``deck`` (metadata) and ``slides`` (an ordered list).

.. code-block:: yaml

    deck:
      title: "Network Enforcement Substrate"
      subtitle: "How 802.1X, SGT and microsegmentation compose the AAN floor"
      series_kicker: "FEDERAL APPLICATION-AWARE NETWORKING"
      date_code: "2026-07-09 14:30 ET"   # CLI --date-code overrides this
      journey:                            # title-slide journey strip labels
        - "Identify"
        - "Authorize"
        - "Segment"
        - "Enforce"
      sources:                            # optional: extends the SRC link-map
        cisco_trustsec: "https://www.cisco.com/c/en/us/products/ios-nx-os-software/trustsec/index.html"

    slides:
      # --- title -----------------------------------------------------------
      - type: title
        notes: "Presenter prose for the opening slide ..."
        vendor_sources: ["cisco_trustsec"]   # optional; resolves against SRC+sources

      # --- section divider -------------------------------------------------
      - type: section
        kicker: "PART ONE"
        title: "Why the wire needs identity"
        callout: "One sentence that frames the section."   # optional banner
        callout_accent: navy                               # navy|teal|amber|red|neutral
        notes: "..."

      # --- content (cards + optional callout) ------------------------------
      - type: content
        kicker: "THE ENFORCEMENT PLANES"
        title: "Four planes carry the policy"
        layout: grid            # grid = 2x2 (up to 4 cards); columns = 3-up
        cards:
          - title: "802.1X"
            body: "Port-based authentication gates the access edge."
            accent: navy
          - title: "SGT"
            body: "Security Group Tags travel with the packet."
            accent: teal
            status: "Enforced"    # 3-up (columns) only: bottom status pill
        callout:                  # optional full-width banner under the cards
          text: "Enforcement is a stack, not a switch."
          accent: navy
        notes: "..."
        vendor_sources: ["https://example.gov/whitepaper"]

      # --- build-sequence (progressive reveal; ONE slide per step) ---------
      - type: build-sequence
        kicker: "THE ENFORCEMENT SEQUENCE"
        title: "How a session is admitted"
        notes: "Fallback notes if a step omits its own."
        steps:
          - label: "Authenticate"
            sub: "802.1X / EAP-TLS"
            accent: navy
            detail: "The supplicant presents a certificate; the switch relays ..."
            notes: "Per-step presenter prose for the slide where this step is active."
          - label: "Authorize"
            sub: "dACL / SGT"
            accent: teal
            detail: "ISE returns an authorization result ..."
            notes: "..."

      # --- table -----------------------------------------------------------
      - type: table
        kicker: "COVERAGE"
        title: "Control-to-plane mapping"
        headers: ["Control", "Plane", "Status"]
        rows:
          - cells: ["AC-4", "Network", "Enforced"]
          - cells: ["SC-7", "Boundary", "Partial"]
            accent: amber          # optional per-row tint
        callout:                   # optional banner under the table
          text: "Boundary enforcement remains the open seam."
          accent: amber
        notes: "..."

      # --- closing (white takeaways slide; never dark) ---------------------
      - type: closing
        kicker: "WHERE THIS LEAVES US"
        title: "The substrate is necessary but not sufficient"
        cards:                     # optional takeaway cards (grid)
          - title: "Necessary"
            body: "No AAN posture without an enforcing floor."
            accent: teal
        callout:
          text: "Next: the evidence fabric that proves it."
          accent: navy
        notes: "..."

Every slide dict carries ``notes`` (required, non-empty) and an optional
``vendor_sources`` list. ``vendor_sources`` entries are either full URLs, a
``SRC`` / ``deck.sources`` key, or a literal ``uiao repo: <path>`` fragment; all
are joined into a single ``Vendor sources: <a> | <b> | ...`` trailer line.

Programmatic entry point: :func:`build_aan_deck`.
"""

from __future__ import annotations

import contextlib
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.text.text import TextFrame
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Canvas
# ---------------------------------------------------------------------------
SLIDE_W_EMU = 9144000  # 10.0 in
SLIDE_H_EMU = 5143500  # 5.625 in

# ---------------------------------------------------------------------------
# Palette (exact hex from OrgComp_DECK_STYLE_NOTES.md)
# ---------------------------------------------------------------------------
NAVY = "1F3A5F"  # primary ink, titles, table header fill
TEAL = "1A9E8F"  # accent / affirmative
AMBER = "D4860B"  # active-step ring, caution
RED = "C0392B"  # failure / urgency
BODY_INK = "1E293B"  # slate near-black body text
MUTED = "5B6B7C"  # secondary text, footers
AMBER_DARK = "8A6210"  # amber text on an amber tint
BORDER = "C9D4E0"  # card / panel border, upcoming-step dot

WHITE = "FFFFFF"  # slide background + cards
NAVY_TINT = "EDF3F9"  # thesis / callout banner
NEUTRAL_PANEL = "F3F5F8"  # neutral panel, upcoming chip
TEAL_TINT = "E6F5F2"
AMBER_TINT = "FBF3E4"  # disclaimer / caution
RED_TINT = "F9EAE8"

# accent name -> (ink hex, tint hex)
ACCENTS: dict[str, tuple[str, str]] = {
    "navy": (NAVY, NAVY_TINT),
    "teal": (TEAL, TEAL_TINT),
    "amber": (AMBER, AMBER_TINT),
    "red": (RED, RED_TINT),
    "neutral": (MUTED, NEUTRAL_PANEL),
}
# rotating palette used for journey-strip dots and default card cycling
ROTATION = ["navy", "teal", "amber", "red", "neutral"]

# Fonts
FONT = "Arial"
FONT_MONO = "Consolas"

# ---------------------------------------------------------------------------
# Source link-map (extendable per-deck via the spec's `sources:` key)
# ---------------------------------------------------------------------------
SRC: dict[str, str] = {
    "uiao_repo": "https://github.com/WhalerMike/uiao",
    "cisco_trustsec": "https://www.cisco.com/c/en/us/products/ios-nx-os-software/trustsec/index.html",
    "cisco_ise": "https://www.cisco.com/c/en/us/products/security/identity-services-engine/index.html",
    "nist_800_53": "https://csrc.nist.gov/pubs/sp/800/53/r5/upd1/final",
    "nist_800_207": "https://csrc.nist.gov/pubs/sp/800/207/final",
    "cisa_ztmm": "https://www.cisa.gov/zero-trust-maturity-model",
    "ieee_802_1x": "https://1.ieee802.org/security/802-1x/",
    "fedramp_20x": "https://www.fedramp.gov/20x/",
}


def resolve_source(entry: str, sources: dict[str, str]) -> str:
    """Resolve a ``vendor_sources`` entry to a display string.

    Full URLs and literal ``uiao repo:`` fragments pass through unchanged; bare
    keys are resolved against the merged ``SRC`` + deck ``sources`` map, and an
    unresolved key is returned verbatim (so authoring typos surface visibly).
    """
    entry = str(entry).strip()
    low = entry.lower()
    if low.startswith(("http://", "https://", "uiao repo:", "repo:")):
        return entry.replace("repo:", "uiao repo:", 1) if low.startswith("repo:") else entry
    return sources.get(entry, entry)


def build_vendor_sources_line(entries: list[str], sources: dict[str, str]) -> str:
    """Render the single ``Vendor sources: a | b | ...`` trailer line."""
    resolved = [resolve_source(e, sources) for e in entries if str(e).strip()]
    return "Vendor sources: " + " | ".join(resolved) if resolved else ""


# ---------------------------------------------------------------------------
# Low-level shape helpers
# ---------------------------------------------------------------------------
def _hex(color: str) -> RGBColor:
    rgb: RGBColor = RGBColor.from_string(color)
    return rgb


def _set_white_background(slide: Slide) -> None:
    """Force a solid white slide background (enforces the no-dark-slides rule)."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _hex(WHITE)


def _no_autofit(tf: TextFrame) -> None:
    # keep author-specified font sizes; python-pptx default is already NONE
    tf.word_wrap = True


def add_textbox(
    slide: Slide,
    left: float,
    top: float,
    width: float,
    height: float,
    paragraphs: list[dict[str, Any]],
    *,
    anchor: Any = MSO_ANCHOR.TOP,
) -> BaseShape:
    """Add a text box.  ``paragraphs`` is a list of paragraph dicts::

        {"runs": [{"text": str, "size": pt, "color": hex, "bold": bool,
                   "italic": bool, "font": name}],
         "align": PP_ALIGN.*, "space_after": pt, "space_before": pt}

    A shorthand single-run paragraph may pass ``text``/``size``/``color``/``bold``
    directly on the paragraph dict.
    """
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    _no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        if para.get("space_after") is not None:
            p.space_after = Pt(para["space_after"])
        if para.get("space_before") is not None:
            p.space_before = Pt(para["space_before"])
        runs = para.get("runs")
        if runs is None:
            runs = [
                {
                    "text": para.get("text", ""),
                    "size": para.get("size", 10),
                    "color": para.get("color", BODY_INK),
                    "bold": para.get("bold", False),
                    "italic": para.get("italic", False),
                    "font": para.get("font", FONT),
                }
            ]
        for rspec in runs:
            r = p.add_run()
            r.text = rspec.get("text", "")
            r.font.size = Pt(rspec.get("size", 10))
            r.font.bold = rspec.get("bold", False)
            r.font.italic = rspec.get("italic", False)
            r.font.name = rspec.get("font", FONT)
            r.font.color.rgb = _hex(rspec.get("color", BODY_INK))
    return box


def add_rounded_rect(
    slide: Slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: str | None = WHITE,
    line: str | None = BORDER,
    line_pt: float = 0.75,
    radius: float = 0.08,
) -> Shape:
    """Add a rounded rectangle with optional fill/border."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    with contextlib.suppress(IndexError, ValueError):
        shp.adjustments[0] = radius
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _hex(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _hex(line)
        shp.line.width = Pt(line_pt)
    shp.shadow.inherit = False
    if shp.has_text_frame:
        shp.text_frame.word_wrap = True
    return shp


def add_dot(slide: Slide, left: float, top: float, size: float, fill: str, *, number: str | None = None) -> BaseShape:
    """Add a colored oval dot, optionally with a white bold number centered."""
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    dot.fill.solid()
    dot.fill.fore_color.rgb = _hex(fill)
    dot.line.fill.background()
    dot.shadow.inherit = False
    tf = dot.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    if number is not None:
        r = p.add_run()
        r.text = str(number)
        r.font.size = Pt(max(8, int(size * 34)))
        r.font.bold = True
        r.font.name = FONT
        r.font.color.rgb = _hex(WHITE)
    return dot


def add_connector(
    slide: Slide, x1: float, y1: float, x2: float, y2: float, color: str, width_pt: float = 1.25
) -> BaseShape:
    """Add a straight connector line."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = _hex(color)
    conn.line.width = Pt(width_pt)
    conn.shadow.inherit = False
    return conn


# ---------------------------------------------------------------------------
# Speaker notes
# ---------------------------------------------------------------------------
def set_notes(slide: Slide, notes: str, vendor_sources: list[str] | None, sources: dict[str, str]) -> None:
    """Attach presenter prose (+ optional Vendor sources line) to a slide."""
    text = (notes or "").rstrip()
    if vendor_sources:
        line = build_vendor_sources_line(vendor_sources, sources)
        if line:
            text = f"{text}\n\n{line}" if text else line
    notes_tf = slide.notes_slide.notes_text_frame
    assert notes_tf is not None
    notes_tf.text = text


# ---------------------------------------------------------------------------
# Standard slide chrome
# ---------------------------------------------------------------------------
def add_chrome(slide: Slide, kicker: str, title: str, footer: str, page_no: int | None) -> None:
    """ALL-CAPS teal kicker, navy title, muted footer + right-aligned page no."""
    if kicker:
        add_textbox(
            slide,
            0.55,
            0.28,
            9.0,
            0.30,
            [{"text": kicker.upper(), "size": 12, "color": TEAL, "bold": True}],
        )
    if title:
        add_textbox(
            slide,
            0.55,
            0.56,
            8.9,
            0.62,
            [{"text": title, "size": 22, "color": NAVY, "bold": True}],
            anchor=MSO_ANCHOR.MIDDLE,
        )
    if footer:
        add_textbox(slide, 0.55, 5.32, 8.0, 0.24, [{"text": footer, "size": 8, "color": MUTED}])
    if page_no is not None:
        add_textbox(
            slide,
            9.20,
            5.32,
            0.35,
            0.24,
            [{"text": str(page_no), "size": 8, "color": MUTED, "align": PP_ALIGN.RIGHT}],
        )


def _new_slide(prs: PresentationType) -> Slide:
    slide: Slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_white_background(slide)
    return slide


# ---------------------------------------------------------------------------
# Title slide
# ---------------------------------------------------------------------------
def render_title_slide(prs: PresentationType, deck: dict[str, Any], date_code: str) -> None:
    slide = _new_slide(prs)
    add_textbox(
        slide,
        0.65,
        0.72,
        8.7,
        0.32,
        [{"text": str(deck.get("series_kicker", "")).upper(), "size": 13, "color": TEAL, "bold": True}],
    )
    add_textbox(
        slide,
        0.65,
        1.08,
        8.7,
        0.85,
        [{"text": deck.get("title", ""), "size": 30, "color": NAVY, "bold": True}],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        0.65,
        1.98,
        8.5,
        0.70,
        [{"text": deck.get("subtitle", ""), "size": 13, "color": MUTED}],
    )

    # journey strip -------------------------------------------------------
    labels = list(deck.get("journey", []) or [])
    if labels:
        n = len(labels)
        total_w = 8.7
        left0 = 0.65
        gap = 0.2
        chip_w = max(1.2, (total_w - (n - 1) * gap) / n)
        chip_h = 0.95
        top = 3.0
        centers_y = top + chip_h / 2.0
        for i, label in enumerate(labels):
            cx = left0 + i * (chip_w + gap)
            add_rounded_rect(slide, cx, top, chip_w, chip_h, fill=WHITE, line=BORDER, line_pt=0.75)
            accent_hex = ACCENTS[ROTATION[i % len(ROTATION)]][0]
            add_dot(slide, cx + chip_w / 2.0 - 0.07, top + 0.14, 0.14, accent_hex)
            add_textbox(
                slide,
                cx + 0.06,
                top + 0.36,
                chip_w - 0.12,
                0.5,
                [{"text": label, "size": 9.5, "color": NAVY, "bold": True, "align": PP_ALIGN.CENTER}],
                anchor=MSO_ANCHOR.MIDDLE,
            )
            if i < n - 1:
                x1 = cx + chip_w
                x2 = cx + chip_w + gap
                add_connector(slide, x1, centers_y, x2, centers_y, MUTED, 1.0)

    # Date Code line ------------------------------------------------------
    add_textbox(
        slide,
        0.65,
        4.52,
        8.7,
        0.26,
        [{"text": f"Date Code: {date_code}", "size": 10, "color": MUTED, "bold": True}],
    )
    # Draft disclaimer ----------------------------------------------------
    disc = add_rounded_rect(slide, 0.65, 4.85, 8.7, 0.52, fill=AMBER_TINT, line=AMBER, line_pt=0.75, radius=0.12)
    tf = disc.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    disclaimer = deck.get(
        "disclaimer",
        "DRAFT PROPOSAL — pre-decisional working material for architecture review. "
        "Not an approved standard, authorization, or agency position.",
    )
    r = p.add_run()
    r.text = disclaimer
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = _hex(AMBER_DARK)

    set_notes(slide, deck.get("_title_notes", ""), deck.get("_title_vendor_sources"), deck["_sources"])


# ---------------------------------------------------------------------------
# Section divider
# ---------------------------------------------------------------------------
def render_section_slide(
    prs: PresentationType, spec: dict[str, Any], page_no: int, footer: str, sources: dict[str, str]
) -> None:
    slide = _new_slide(prs)
    add_chrome(slide, spec.get("kicker", ""), spec.get("title", ""), footer, page_no)
    callout = spec.get("callout")
    if callout:
        _render_callout(slide, callout, spec.get("callout_accent", "navy"), top=2.4, height=0.9, size=13)
    set_notes(slide, spec.get("notes", ""), spec.get("vendor_sources"), sources)


# ---------------------------------------------------------------------------
# Callout / thesis banner
# ---------------------------------------------------------------------------
def _render_callout(
    slide: Slide,
    callout: str | dict[str, Any],
    default_accent: str,
    *,
    top: float,
    height: float,
    size: float = 11,
) -> None:
    if isinstance(callout, str):
        text, accent = callout, default_accent
    else:
        text = callout.get("text", "")
        accent = callout.get("accent", default_accent)
    ink, tint = ACCENTS.get(accent, ACCENTS["navy"])
    banner = add_rounded_rect(slide, 0.55, top, 9.0, height, fill=tint, line=BORDER, line_pt=0.75, radius=0.06)
    tf = banner.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = _hex(ink)


# ---------------------------------------------------------------------------
# Content cards
# ---------------------------------------------------------------------------
def _render_card(
    slide: Slide,
    left: float,
    top: float,
    width: float,
    height: float,
    index: int,
    card: dict[str, Any],
    *,
    with_status: bool,
) -> None:
    accent = card.get("accent") or ROTATION[(index - 1) % len(ROTATION)]
    ink, tint = ACCENTS.get(accent, ACCENTS["navy"])
    fill = card.get("fill", WHITE)
    add_rounded_rect(slide, left, top, width, height, fill=fill, line=BORDER, line_pt=0.75, radius=0.06)
    # numbered dot
    add_dot(slide, left + 0.18, top + 0.16, 0.32, ink, number=str(index))
    # title (to the right of the dot)
    add_textbox(
        slide,
        left + 0.6,
        top + 0.14,
        width - 0.75,
        0.36,
        [{"text": card.get("title", ""), "size": 12, "color": NAVY, "bold": True}],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # body
    body_h = height - (0.9 if with_status else 0.68)
    add_textbox(
        slide,
        left + 0.18,
        top + 0.58,
        width - 0.36,
        body_h,
        [{"text": card.get("body", ""), "size": 9, "color": BODY_INK}],
    )
    # status pill (3-up columns only)
    if with_status and card.get("status"):
        pill = add_rounded_rect(
            slide,
            left + 0.18,
            top + height - 0.46,
            width - 0.36,
            0.32,
            fill=tint,
            line=None,
            radius=0.3,
        )
        tf = pill.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(card["status"])
        r.font.size = Pt(8.5)
        r.font.bold = True
        r.font.name = FONT
        r.font.color.rgb = _hex(ink if accent != "amber" else AMBER_DARK)


def _render_cards(slide: Slide, cards: list[dict[str, Any]], layout: str, *, has_callout: bool) -> None:
    if layout == "columns":
        n = min(len(cards), 3)
        card_w = 2.94
        gap = (9.0 - 3 * card_w) / 2.0  # ~0.09
        top = 1.5
        height = 2.7 if not has_callout else 2.45
        for i in range(n):
            left = 0.55 + i * (card_w + gap)
            _render_card(slide, left, top, card_w, height, i + 1, cards[i], with_status=True)
    else:  # grid 2x2
        card_w = 4.42
        gap_x = 9.0 - 2 * card_w  # 0.16
        if has_callout:
            card_h, t1, t2 = 1.4, 1.45, 3.0
        else:
            card_h, t1, t2 = 1.52, 1.5, 3.25
        positions = [
            (0.55, t1),
            (0.55 + card_w + gap_x, t1),
            (0.55, t2),
            (0.55 + card_w + gap_x, t2),
        ]
        for i, card in enumerate(cards[:4]):
            left, top = positions[i]
            _render_card(slide, left, top, card_w, card_h, i + 1, card, with_status=False)


def render_content_slide(
    prs: PresentationType, spec: dict[str, Any], page_no: int, footer: str, sources: dict[str, str]
) -> None:
    slide = _new_slide(prs)
    add_chrome(slide, spec.get("kicker", ""), spec.get("title", ""), footer, page_no)
    cards = spec.get("cards", []) or []
    layout = spec.get("layout", "grid")
    callout = spec.get("callout")
    if cards:
        _render_cards(slide, cards, layout, has_callout=bool(callout))
    if callout:
        top = 4.55 if layout != "columns" else 4.35
        _render_callout(slide, callout, spec.get("callout_accent", "navy"), top=top, height=0.6)
    set_notes(slide, spec.get("notes", ""), spec.get("vendor_sources"), sources)


# ---------------------------------------------------------------------------
# Build-sequence (progressive reveal) — one slide per active step
# ---------------------------------------------------------------------------
def render_build_sequence(
    prs: PresentationType, spec: dict[str, Any], start_page: int, footer: str, sources: dict[str, str]
) -> int:
    """Emit one slide per step; returns the next page number."""
    steps = spec.get("steps", []) or []
    n = len(steps)
    kicker = spec.get("kicker", "")
    title = spec.get("title", "")
    fallback_notes = spec.get("notes", "")

    total_w = 9.0
    left0 = 0.55
    gap = 0.32
    chip_w = (total_w - (n - 1) * gap) / n if n else total_w
    chip_h = 0.85
    strip_top = 1.42
    conn_y = strip_top + chip_h / 2.0

    page = start_page
    for active in range(n):
        slide = _new_slide(prs)
        add_chrome(slide, kicker, title, footer, page)

        # connectors first (behind chips)
        for i in range(n - 1):
            x1 = left0 + i * (chip_w + gap) + chip_w
            x2 = left0 + (i + 1) * (chip_w + gap)
            color = BODY_INK if i < active else BORDER
            width = 1.75 if i < active else 1.25
            add_connector(slide, x1, conn_y, x2, conn_y, color, width)

        # chips
        for i, step in enumerate(steps):
            cx = left0 + i * (chip_w + gap)
            accent = step.get("accent") or ROTATION[i % len(ROTATION)]
            ink, tint = ACCENTS.get(accent, ACCENTS["navy"])
            if i == active:
                # amber ring behind the active chip (inflated 0.06 per side)
                add_rounded_rect(
                    slide,
                    cx - 0.06,
                    strip_top - 0.06,
                    chip_w + 0.12,
                    chip_h + 0.12,
                    fill=WHITE,
                    line=AMBER,
                    line_pt=2.5,
                    radius=0.08,
                )
                chip_fill, dot_fill, text_color = WHITE, ink, NAVY
            elif i < active:  # completed
                chip_fill, dot_fill, text_color = WHITE, ink, NAVY
            else:  # upcoming
                chip_fill, dot_fill, text_color = NEUTRAL_PANEL, BORDER, MUTED
            add_rounded_rect(slide, cx, strip_top, chip_w, chip_h, fill=chip_fill, line=BORDER, line_pt=0.75)
            add_dot(slide, cx + 0.14, strip_top + 0.14, 0.14, dot_fill)
            paras: list[dict[str, Any]] = [
                {
                    "text": f"{i + 1} — {step.get('label', '')}",
                    "size": 9.5,
                    "color": text_color,
                    "bold": True,
                    "space_after": 1,
                }
            ]
            if step.get("sub"):
                paras.append({"text": step["sub"], "size": 8, "color": MUTED if i != active else MUTED})
            add_textbox(slide, cx + 0.32, strip_top + 0.12, chip_w - 0.42, chip_h - 0.2, paras)

        # detail panel (deepens per step) — active step's accent tint
        active_accent = steps[active].get("accent") or ROTATION[active % len(ROTATION)]
        _, panel_tint = ACCENTS.get(active_accent, ACCENTS["navy"])
        add_rounded_rect(slide, 0.55, 2.55, 9.0, 2.6, fill=panel_tint, line=BORDER, line_pt=0.75, radius=0.05)
        detail = steps[active].get("detail", "")
        step_label = steps[active].get("label", "")
        add_textbox(
            slide,
            0.85,
            2.75,
            8.45,
            2.25,
            [
                {"text": f"{active + 1}. {step_label}", "size": 13, "color": NAVY, "bold": True, "space_after": 6},
                {"text": detail, "size": 11, "color": BODY_INK},
            ],
        )

        step_notes = steps[active].get("notes") or fallback_notes
        step_vs = steps[active].get("vendor_sources") or spec.get("vendor_sources")
        set_notes(slide, step_notes, step_vs, sources)
        page += 1
    return page


# ---------------------------------------------------------------------------
# Native styled table
# ---------------------------------------------------------------------------
def render_table_slide(
    prs: PresentationType, spec: dict[str, Any], page_no: int, footer: str, sources: dict[str, str]
) -> None:
    slide = _new_slide(prs)
    add_chrome(slide, spec.get("kicker", ""), spec.get("title", ""), footer, page_no)

    headers = spec.get("headers", []) or []
    rows = spec.get("rows", []) or []
    n_cols = len(headers)
    n_rows = len(rows) + 1
    callout = spec.get("callout")
    table_h = 3.0 if callout else 3.7
    gfx = slide.shapes.add_table(n_rows, n_cols, Inches(0.55), Inches(1.4), Inches(9.0), Inches(table_h))
    table = gfx.table

    # header row
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _hex(NAVY)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = str(h)
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.name = FONT
        r.font.color.rgb = _hex(WHITE)

    # body rows
    for ridx, row in enumerate(rows, start=1):
        cells = row.get("cells", []) if isinstance(row, dict) else list(row)
        accent = row.get("accent") if isinstance(row, dict) else None
        row_fill = ACCENTS[accent][1] if accent in ACCENTS else WHITE
        for c in range(n_cols):
            cell = table.cell(ridx, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex(row_fill)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = str(cells[c]) if c < len(cells) else ""
            r.font.size = Pt(9)
            r.font.name = FONT
            if c == 0:
                r.font.bold = True
                r.font.color.rgb = _hex(NAVY)
            else:
                r.font.color.rgb = _hex(BODY_INK)

    if callout:
        _render_callout(slide, callout, spec.get("callout_accent", "navy"), top=1.4 + table_h + 0.15, height=0.55)
    set_notes(slide, spec.get("notes", ""), spec.get("vendor_sources"), sources)


# ---------------------------------------------------------------------------
# Closing (white takeaways) slide
# ---------------------------------------------------------------------------
def render_closing_slide(
    prs: PresentationType, spec: dict[str, Any], page_no: int, footer: str, sources: dict[str, str]
) -> None:
    # a closing slide is just a white content slide; enforced non-dark
    render_content_slide(prs, spec, page_no, footer, sources)


# ---------------------------------------------------------------------------
# Orchestrator
# ---------------------------------------------------------------------------
_RENDERERS = {
    "section": render_section_slide,
    "content": render_content_slide,
    "table": render_table_slide,
    "closing": render_closing_slide,
}


def build_presentation(spec: dict[str, Any], date_code: str) -> PresentationType:
    """Build a :class:`Presentation` from a parsed deck spec + date code."""
    deck = dict(spec.get("deck", {}) or {})
    slides = spec.get("slides", []) or []

    # merged source map: SRC <- deck.sources
    sources = dict(SRC)
    sources.update(deck.get("sources", {}) or {})
    deck["_sources"] = sources

    footer = deck.get("footer", f"{deck.get('title', 'AAN Briefing')} · Date Code: {date_code}")

    prs: PresentationType = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)

    page = 1
    for slide_spec in slides:
        stype = slide_spec.get("type")
        if stype == "title":
            # title slide carries its own notes fields; stash onto deck
            deck["_title_notes"] = slide_spec.get("notes", "")
            deck["_title_vendor_sources"] = slide_spec.get("vendor_sources")
            render_title_slide(prs, deck, date_code)
            page += 1
        elif stype == "build-sequence":
            page = render_build_sequence(prs, slide_spec, page, footer, sources)
        elif stype in _RENDERERS:
            _RENDERERS[stype](prs, slide_spec, page, footer, sources)
            page += 1
        else:
            raise ValueError(f"Unknown slide type: {stype!r}")

    return prs


def build_aan_deck(
    *,
    spec_path: str | Path,
    out_path: str | Path,
    date_code: str,
) -> Path:
    """Load a deck-spec YAML, build the deck, save to ``out_path``, return it.

    Parameters
    ----------
    spec_path:
        Path to the deck-spec YAML (schema documented in the module docstring).
    out_path:
        Destination ``.pptx`` path (parent dirs are created).
    date_code:
        ``YYYY-MM-DD HH:MM ET`` string rendered on the title slide. Passed
        explicitly by the caller; never derived from the machine clock/``TZ``.
    """
    spec_path = Path(spec_path)
    out_path = Path(out_path)
    with spec_path.open("r", encoding="utf-8") as fh:
        spec = yaml.safe_load(fh)
    if not isinstance(spec, dict):
        raise ValueError(f"Deck spec {spec_path} did not parse to a mapping")

    # allow the spec to supply date_code when the caller passes an empty string
    effective = date_code or str(spec.get("deck", {}).get("date_code", "")).strip()
    if not effective:
        raise ValueError("No date_code provided (pass --date-code or set deck.date_code in the spec)")

    prs = build_presentation(spec, effective)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path
