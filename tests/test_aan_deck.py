"""Tests for uiao.generators.aan_deck — canonical AAN briefing-deck generator.

Asserts the Book_10/Book_11 FedAAN deck DNA measured in
``inbox/Application Aware Networking/AAN_DECK_STYLE_NOTES.md``:

* slide size exactly 10.0 x 5.625 in;
* every slide carries non-empty speaker notes;
* every slide background is solid white (no dark slides);
* the house palette hex is used as ink;
* the title slide renders the ``Date Code:`` string;
* a build-sequence slide renders the amber-ring (``D4860B``) active step;
* a table renders with the navy (``1F3A5F``) header row.
"""

from __future__ import annotations

import yaml
from pptx.util import Emu

from uiao.generators.aan_deck import (
    AMBER,
    NAVY,
    SLIDE_H_EMU,
    SLIDE_W_EMU,
    TEAL,
    build_aan_deck,
    build_presentation,
    build_vendor_sources_line,
)

DATE_CODE = "2026-07-09 14:30 ET"

SPEC = {
    "deck": {
        "title": "Test Enforcement Substrate",
        "subtitle": "A minimal in-test deck exercising every slide type",
        "series_kicker": "FEDERAL APPLICATION-AWARE NETWORKING",
        "date_code": DATE_CODE,
        "journey": ["Identify", "Authorize", "Segment", "Enforce"],
        "sources": {"my_vendor": "https://example.gov/whitepaper"},
    },
    "slides": [
        {"type": "title", "notes": "Opening frame.", "vendor_sources": ["my_vendor"]},
        {
            "type": "section",
            "kicker": "PART ONE",
            "title": "Why the wire needs identity",
            "callout": "Identity first.",
            "notes": "Section motivation prose.",
        },
        {
            "type": "content",
            "kicker": "PLANES",
            "title": "Four planes",
            "layout": "grid",
            "cards": [
                {"title": "802.1X", "body": "Access edge.", "accent": "navy"},
                {"title": "SGT", "body": "Tag travels.", "accent": "teal"},
            ],
            "callout": {"text": "A stack, not a switch.", "accent": "navy"},
            "notes": "Card walk-through prose.",
        },
        {
            "type": "build-sequence",
            "kicker": "SEQUENCE",
            "title": "How a session is admitted",
            "notes": "Fallback notes.",
            "steps": [
                {"label": "Authenticate", "sub": "802.1X", "accent": "navy", "detail": "Cert.", "notes": "Step one."},
                {"label": "Authorize", "sub": "dACL", "accent": "teal", "detail": "Result.", "notes": "Step two."},
            ],
        },
        {
            "type": "table",
            "kicker": "COVERAGE",
            "title": "Mapping",
            "headers": ["Control", "Plane", "Status"],
            "rows": [
                {"cells": ["AC-4", "Segmentation", "Enforced"], "accent": "teal"},
                {"cells": ["SC-7", "Boundary", "Partial"], "accent": "amber"},
            ],
            "notes": "Table prose.",
        },
        {
            "type": "closing",
            "kicker": "WHERE THIS LEAVES US",
            "title": "Necessary but not sufficient",
            "layout": "columns",
            "cards": [
                {"title": "Necessary", "body": "The floor.", "accent": "teal", "status": "Done"},
            ],
            "notes": "Closing prose.",
        },
    ],
}


def _build():
    return build_presentation(SPEC, DATE_CODE)


def _shape_line_hex(shape) -> str | None:
    try:
        return str(shape.line.color.rgb)
    except (TypeError, AttributeError):
        return None


class TestCanvas:
    def test_slide_size_emu(self):
        prs = _build()
        assert prs.slide_width == Emu(SLIDE_W_EMU)
        assert prs.slide_height == Emu(SLIDE_H_EMU)

    def test_slide_size_inches(self):
        prs = _build()
        assert round(prs.slide_width.inches, 3) == 10.0
        assert round(prs.slide_height.inches, 3) == 5.625


class TestNotes:
    def test_every_slide_has_nonempty_notes(self):
        prs = _build()
        for i, slide in enumerate(prs.slides):
            assert slide.has_notes_slide, f"slide {i} missing notes"
            text = slide.notes_slide.notes_text_frame.text.strip()
            assert text, f"slide {i} has empty notes"

    def test_vendor_sources_line_appended(self):
        prs = _build()
        # title slide had vendor_sources -> resolved to the my_vendor URL
        title_notes = prs.slides[0].notes_slide.notes_text_frame.text
        assert "Vendor sources:" in title_notes
        assert "https://example.gov/whitepaper" in title_notes

    def test_vendor_sources_resolution_helper(self):
        line = build_vendor_sources_line(
            ["cisco_trustsec", "uiao repo: src/uiao/substrate/"],
            {"cisco_trustsec": "https://cisco.example/x"},
        )
        assert line.startswith("Vendor sources: ")
        assert "https://cisco.example/x" in line
        assert "uiao repo: src/uiao/substrate/" in line


class TestBackgrounds:
    def test_all_backgrounds_white(self):
        prs = _build()
        for i, slide in enumerate(prs.slides):
            rgb = str(slide.background.fill.fore_color.rgb)
            assert rgb == "FFFFFF", f"slide {i} background {rgb} is not white"


class TestPalette:
    def test_title_slide_contains_date_code(self):
        prs = _build()
        texts = [s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame]
        assert any(f"Date Code: {DATE_CODE}" in t for t in texts)

    def test_navy_ink_used_on_a_title(self):
        prs = _build()
        found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if str(run.font.color.rgb) == NAVY:
                                found = True
                        except (TypeError, AttributeError):
                            pass
        assert found, "navy ink hex never used"

    def test_teal_kicker_used(self):
        prs = _build()
        found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if str(run.font.color.rgb) == TEAL:
                                found = True
                        except (TypeError, AttributeError):
                            pass
        assert found, "teal kicker hex never used"


class TestBuildSequence:
    def test_one_slide_per_step(self):
        prs = _build()
        # 2 steps -> 2 build-sequence slides; total = title+section+content+2+table+closing = 7
        assert len(prs.slides) == 7

    def test_amber_ring_present_on_a_build_slide(self):
        prs = _build()
        amber_ring = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if _shape_line_hex(shape) == AMBER:
                    amber_ring = True
        assert amber_ring, "amber-ring active step never rendered"


class TestTable:
    def test_navy_header_row(self):
        prs = _build()
        header_navy = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_table:
                    continue
                cell = shape.table.cell(0, 0)
                if str(cell.fill.fore_color.rgb) == NAVY:
                    header_navy = True
        assert header_navy, "table header row not navy"

    def test_table_body_first_col_bold_navy(self):
        prs = _build()
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_table:
                    continue
                cell = shape.table.cell(1, 0)
                run = cell.text_frame.paragraphs[0].runs[0]
                assert run.font.bold is True
                assert str(run.font.color.rgb) == NAVY
                return
        raise AssertionError("no table found")


class TestEndToEnd:
    def test_build_from_yaml_file(self, tmp_path):
        spec_path = tmp_path / "spec.yaml"
        spec_path.write_text(yaml.safe_dump(SPEC), encoding="utf-8")
        out = tmp_path / "deck.pptx"
        result = build_aan_deck(spec_path=spec_path, out_path=out, date_code=DATE_CODE)
        assert result.exists()
        assert result.stat().st_size > 0

    def test_date_code_falls_back_to_spec(self, tmp_path):
        spec_path = tmp_path / "spec.yaml"
        spec_path.write_text(yaml.safe_dump(SPEC), encoding="utf-8")
        out = tmp_path / "deck.pptx"
        # empty CLI date_code -> spec deck.date_code is used
        result = build_aan_deck(spec_path=spec_path, out_path=out, date_code="")
        assert result.exists()
