#!/usr/bin/env python3
"""Gate 4 — tracked derivatives must carry their source's Date Code.

Authoring-spec §1: one Date Code per book, bumped on every content change,
rendered immediately so source and derivatives match. This checks, for every
GIT-TRACKED derivative:

  - book .docx  vs its .qmd frontmatter `date:`
  - deck .pptx  vs its decks/<name>.yaml `date_code:`
  - deck yaml   vs its book qmd `date:` (spec: deck and book stay in sync)

Untracked derivatives are ignored (they are local build products). A tracked
derivative with no extractable Date Code is an error — it can never be proven
fresh.

Usage:
    python check_derivative_freshness.py        # exit 1 on any mismatch
"""

from __future__ import annotations

import re
import subprocess
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
DC_RE = re.compile(r"Date Code:\s*(\d{4}-\d{2}-\d{2} \d{2}:\d{2} ET)")
QMD_DATE_RE = re.compile(r'^date:\s*"([^"]+)"', re.M)
YAML_DC_RE = re.compile(r'date_code:\s*"([^"]+)"')


def tracked(ext: str) -> list[Path]:
    out = subprocess.run(["git", "ls-files", f"*{ext}"], capture_output=True, text=True, cwd=HERE).stdout.splitlines()
    return [
        HERE / p for p in out if (HERE / p).exists() and "/" not in p.replace("\\", "/") or True and (HERE / p).exists()
    ]


def docx_code(p: Path) -> str | None:
    from docx import Document

    d = Document(str(p))
    for para in d.paragraphs:
        m = DC_RE.search(para.text)
        if m:
            return m.group(1)
    return None


def pptx_code(p: Path) -> str | None:
    from pptx import Presentation

    prs = Presentation(str(p))
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                m = DC_RE.search(sh.text_frame.text)
                if m:
                    return m.group(1)
    return None


def main() -> int:
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8")
    errors: list[str] = []
    checked = 0

    for docx in tracked(".docx"):
        qmd = docx.with_suffix(".qmd")
        if not qmd.exists():
            continue
        src = QMD_DATE_RE.search(qmd.read_text(encoding="utf-8", errors="replace"))
        if not src:
            continue
        checked += 1
        code = docx_code(docx)
        if code is None:
            errors.append(f"{docx.name}: tracked docx carries no extractable Date Code")
        elif code != src.group(1):
            errors.append(f"{docx.name}: docx Date Code {code} != qmd date {src.group(1)} — re-render")

    for pptx in tracked(".pptx"):
        base = pptx.stem
        # deck spec name convention: Vol_IX_Book_05 -> decks/Vol_IX_Book_05.yaml
        m = re.match(r"(Vol_[0IVX]+_Book_\w+?)_OrgComp", base)
        spec = HERE / "decks" / (m.group(1) + ".yaml") if m else None
        if not spec or not spec.exists():
            continue
        sd = YAML_DC_RE.search(spec.read_text(encoding="utf-8", errors="replace"))
        if not sd:
            continue
        checked += 1
        code = pptx_code(pptx)
        if code is None:
            errors.append(f"{pptx.name}: tracked pptx carries no extractable Date Code")
        elif code != sd.group(1):
            errors.append(f"{pptx.name}: pptx Date Code {code} != deck spec {sd.group(1)} — regenerate")
        # deck spec vs the book qmd
        qmd = next(iter(HERE.glob(m.group(1) + "_OrgComp_*.qmd")), None)
        if qmd:
            qd = QMD_DATE_RE.search(qmd.read_text(encoding="utf-8", errors="replace"))
            if qd and qd.group(1) != sd.group(1):
                errors.append(
                    f"{spec.name}: deck date_code {sd.group(1)} != book qmd date {qd.group(1)} — sync per spec §1"
                )

    print("AAN derivative-freshness check")
    print("=" * 44)
    print(f"Tracked derivative pairs checked: {checked}")
    if errors:
        print("\nERRORS:")
        for e in errors:
            print(f"  - {e}")
        return 1
    print("\nOK — every tracked derivative carries its source's Date Code.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
